from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.ai.generation.llm_client import LLMClient
from app.core.logging import logger

CHUNK_SYSTEM_PROMPT = """Bạn là chuyên gia trích xuất và tổ chức nội dung học thuật.
Nhiệm vụ: TRICH XUAT va TO CHUC LAI thong tin tu CHUNK DUOC CUNG CAP thanh slide.

QUY TAC BAT BUOC:
1) CHI TRICH XUAT - KHONG DUOC TAO MOI hoac THEM kien thuc ngoai chunk.
2) Moi bullet phai la CAU TRUC XUAT TRUC TIEP tu chunk (facts, definitions, steps, examples CO SAN).
3) Moi bullet phai dai 10-20 tu, BAO TOM day du y nghia cua doan text goc.
4) NEU chunk co danh sach (list), dinh nghia, or buoc thuc hien → GIAI THICH RO RANG trong bullet.
5) Moi slide phai co 3-5 bullets CHI TIET.
6) Neu chunk ngan hoac it thong tin, van phai TRICH XUAT het thong tin co duoc.

VI DU CACH TRICH XUAT:
- Chunk: "Machine learning uses data to build models. It does not require explicit programming."
  → Bullet: "Machine Learning la ky thuat su dung du lieu de xay dung mo hinh du doan, khong can lap trinh logic ro rang"

- Chunk: "K-means divides data into K clusters. Each cluster has a centroid."
  → Bullet: "Thuat toan K-means chia du lieu thanh K nhom (clusters), moi nhom co mot diem trung tam (centroid)"

CAM:
- Them y kien ca nhan
- Tong quat hoa qua muc ("Machine Learning rat quan trong")
- Bullets qua ngan khong ro nghia

CHI TRA VE JSON HOP LE, KHONG markdown, KHONG giai thich.
"""

CHUNK_PROMPT_TEMPLATE = """Ten chunk: {chunk_title}
Tone: {tone}
So slide toi da cho chunk nay: {slide_budget}

NOI DUNG CHUNK:
{chunk_content}

METADATA ANH TU TAI LIEU (neu co):
{available_images_json}

YEU CAU OUTPUT CHINH XAC THEO FORMAT SAU:
{{
  "title": "...",
  "slides": [
    {{
      "title": "...",
      "bullets": ["...", "...", "..."],
      "image_source_id": "img_0"
    }}
  ]
}}

RANG BUOC:
- So bullets moi slide: 3-5
- Bullet phai TRICH XUAT TRUC TIEP tu chunk (10-20 tu/bullet, day du y nghia)
- image_source_id: CHI dung image_id tu METADATA ANH ben tren, KHONG tao moi
  + Neu co anh phu hop → dung image_id do (vi du: "img_0", "img_1")
  + Neu khong co anh phu hop → DE TRONG ("image_source_id": null)
- KHONG duoc them hoac bien doan noi dung chunk
"""

REFINE_SYSTEM_PROMPT = """Ban la reviewer toi uu bo slide hoc tap.
Nhiem vu: CHI sap xep lai va loai bo trung lap - GIU NGUYEN TOI DA noi dung goc.

QUY TAC:
1) Remove slides trung lap (tieu de va bullets giong >70%).
2) Sap xep trinh tu logic (tong quat → chi tiet, co ban → nang cao).
3) GIU NGUYEN 100% bullets da tot - CHI loai bo neu trung lap hoan toan.
4) Dam bao moi slide co 3-5 bullets.
5) NEU phai cat bot de <= gioi han: uu tien GIU slides co nhieu thong tin cu the nhat.
6) KHONG DUOC sua doi, rut gon, hoac viet lai bullets da co - chi sap xep va chon loc.

CHI TRA VE JSON HOP LE (khong markdown):
{
  "title": "...",
  "slides": [
    {
      "title": "...",
      "bullets": ["...", "...", "..."],
      "image_source_id": "img_0"
    }
  ]
}
"""


BRAND_PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
BRAND_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
HIGHLIGHT_BG = RGBColor(0xFE, 0xF3, 0xC7)
HIGHLIGHT_TEXT = RGBColor(0x92, 0x40, 0x0E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1E, 0x1E, 0x2E)
BODY_TEXT = RGBColor(0x4B, 0x55, 0x63)
SUBTITLE_TEXT = RGBColor(0x6B, 0x72, 0x80)
IMG_PLACEHOLDER_BG = RGBColor(0xE0, 0xE7, 0xFF)


class SlideGenerator:
    def __init__(self) -> None:
        self.llm = LLMClient()

    def generate_from_chunk(
        self,
        chunk: dict,
        tone: str = "teacher",
        slide_budget: int = 2,
        available_images: list[dict] | None = None,
    ) -> list[dict]:
        fallback = {
            "title": chunk.get("title", "Nội dung"),
            "slides": [
                {
                    "title": chunk.get("title", "Nội dung chính"),
                    "bullets": self._fallback_bullets_from_chunk(chunk.get("content", "")),
                }
            ],
        }

        # Truncate chunk content to 8000 chars max to reduce token cost
        # (avg 1 char ~= 0.3 tokens → 8000 chars ~= 2400 tokens)
        chunk_content = (chunk.get("content", "") or "")[:8000]

        prompt = CHUNK_PROMPT_TEMPLATE.format(
            chunk_title=chunk.get("title", "Nội dung"),
            tone=tone,
            slide_budget=max(1, min(2, slide_budget)),
            chunk_content=chunk_content,
            available_images_json=json.dumps(
                {"available_images": available_images or []},
                ensure_ascii=False,
                indent=2,
            ),
        )

        result = self.llm.json_response(CHUNK_SYSTEM_PROMPT, prompt, fallback)
        raw_slides = result.get("slides", fallback["slides"])
        normalized = self._normalize_generated_slides(raw_slides, slide_budget)
        return normalized

    def refine_slides(
        self,
        slides: list[dict],
        presentation_title: str,
        max_slides: int,
    ) -> dict:
        fallback = {
            "title": presentation_title,
            "slides": self._normalize_generated_slides(slides, max_slides),
        }

        user_prompt = (
            f"Title: {presentation_title}\n"
            f"Max slides: {max_slides}\n\n"
            "Slides draft to refine:\n"
            f"{json.dumps(slides, ensure_ascii=False, indent=2)}"
        )

        result = self.llm.json_response(REFINE_SYSTEM_PROMPT, user_prompt, fallback)
        refined_slides = self._normalize_generated_slides(result.get("slides", []), max_slides)

        # Keep output backward-compatible with existing frontend/export code.
        final_slides = [self._to_legacy_slide_model(slide, idx) for idx, slide in enumerate(refined_slides)]
        return {"title": result.get("title") or presentation_title, "slides": final_slides}

    def _normalize_generated_slides(self, slides: list[dict], max_slides: int) -> list[dict]:
        normalized: list[dict] = []
        for slide in slides:
            if not isinstance(slide, dict):
                continue
            title = (slide.get("title") or "Nội dung").strip()
            bullets = slide.get("bullets") or []
            if not isinstance(bullets, list):
                bullets = []

            cleaned_bullets = [self._clean_line(str(item)) for item in bullets if self._clean_line(str(item))]
            cleaned_bullets = [b for b in cleaned_bullets if not self._is_generic_bullet(b)]

            if len(cleaned_bullets) < 3:
                cleaned_bullets = self._pad_bullets(cleaned_bullets, title)

            normalized.append(
                {
                    "title": title,
                    "bullets": cleaned_bullets[:5],
                    "image_source_id": slide.get("image_source_id"),
                    "doc_image_index": slide.get("doc_image_index"),
                }
            )

            if len(normalized) >= max_slides:
                break

        return normalized

    @staticmethod
    def _clean_line(text: str) -> str:
        return re.sub(r"\s+", " ", (text or "").strip())

    @staticmethod
    def _is_generic_bullet(text: str) -> bool:
        """Filter bullets that are too short or vague."""
        t = (text or "").strip()

        # Filter bullets that are too short (< 8 words)
        word_count = len(t.split())
        if word_count < 8:
            return True

        # Filter bullets that are ONLY generic phrases without details
        t_lower = t.lower()
        generic_only_patterns = [
            r"^(gioi thieu|giới thiệu)\s*(về|ve)?\s*$",
            r"^(tong quan|tổng quan)\s*$",
            r"^(ung dung|ứng dụng)\s*(thuc te|thực tiễn)?\s*$",
            r"^(vi du|ví dụ)\s*(minh hoa|minh họa)?\s*$",
            r"^(khai niem|khái niệm)\s*(co ban|cơ bản)?\s*$",
            r"^(dinh nghia|định nghĩa)\s*(ro rang|rõ ràng)?\s*$",
        ]

        import re as regex
        for pattern in generic_only_patterns:
            if regex.match(pattern, t_lower):
                return True

        return False

    def _pad_bullets(self, bullets: list[str], title: str) -> list[str]:
        result = list(bullets)
        defaults = [
            f"Ý chính 1 của phần {title}",
            f"Ý chính 2 của phần {title}",
            f"Chi tiết quan trọng trong phần {title}",
        ]
        for item in defaults:
            if len(result) >= 3:
                break
            result.append(item)
        return result

    def _fallback_bullets_from_chunk(self, chunk_content: str) -> list[str]:
        lines = [self._clean_line(line) for line in re.split(r"[\.\n]", chunk_content) if self._clean_line(line)]
        lines = [line for line in lines if len(line.split()) >= 4]
        if len(lines) >= 3:
            return lines[:5]
        return self._pad_bullets(lines, "nội dung")

    def _to_legacy_slide_model(self, slide: dict, idx: int) -> dict:
        layout = "title_only" if idx == 0 else "text_image"
        elements = [
            {"type": "bullet", "content": slide.get("bullets", [])},
        ]

        image_source_id = slide.get("image_source_id")

        # Only use document images - no external image queries
        if image_source_id:
            elements.append(
                {
                    "type": "image_source",
                    "image_id": image_source_id,
                    "image_context": "Ảnh từ tài liệu nguồn",
                }
            )

        return {
            "title": slide.get("title", "Nội dung"),
            "layout": layout,
            "elements": elements,
            # Fields for compatibility
            "bullets": slide.get("bullets", []),
            "image_source_id": image_source_id,
            "doc_image_index": slide.get("doc_image_index"),
        }

    def export_pptx(
        self,
        outline: dict,
        output_path: str,
        image_map: dict[str, str | None] | None = None,
        doc_image_paths: list[str] | None = None,
    ) -> str:
        image_map = image_map or {}
        doc_image_paths = doc_image_paths or []

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slides_data = outline.get("slides", [])
        presentation_title = outline.get("title", "Bài giảng")

        for idx, slide_data in enumerate(slides_data):
            layout = slide_data.get("layout") or ("title_only" if idx == 0 else "text_image")
            title = slide_data.get("title", f"Slide {idx + 1}")

            bullets, doc_image_index = self._extract_slide_parts(slide_data)

            image_path: str | None = None
            if doc_image_index is not None and 0 <= doc_image_index < len(doc_image_paths):
                candidate = doc_image_paths[doc_image_index]
                if candidate and Path(candidate).exists():
                    image_path = candidate

            if layout == "title_only":
                self._add_title_slide(prs, title, presentation_title, bullets)
            elif layout == "text_only":
                self._add_text_only_slide(prs, title, bullets, idx, len(slides_data))
            else:
                self._add_text_image_slide(
                    prs,
                    title,
                    bullets,
                    image_path,
                    idx,
                    len(slides_data),
                )

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output)
        return str(output)

    def _extract_slide_parts(self, slide_data: dict) -> tuple[list[str], int | None]:
        bullets = list(slide_data.get("bullets") or [])
        doc_image_index = slide_data.get("doc_image_index")

        image_source_id = slide_data.get("image_source_id")
        if doc_image_index is None and isinstance(image_source_id, str) and image_source_id.startswith("img_"):
            try:
                doc_image_index = int(image_source_id.replace("img_", ""))
            except ValueError:
                doc_image_index = None

        for elem in slide_data.get("elements", []):
            etype = elem.get("type")
            if etype == "bullet" and not bullets:
                bullets = list(elem.get("content", []))
            elif etype in ("image_source", "doc_image") and doc_image_index is None:
                eid = elem.get("image_id")
                if isinstance(eid, str) and eid.startswith("img_"):
                    try:
                        doc_image_index = int(eid.replace("img_", ""))
                    except ValueError:
                        pass

        return bullets, doc_image_index

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str, bullets: list[str]) -> None:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BRAND_PRIMARY

        tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.5)).text_frame
        tx.word_wrap = True
        p = tx.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        tx2 = slide.shapes.add_textbox(Inches(2.5), Inches(4.0), Inches(8), Inches(1.0)).text_frame
        tx2.word_wrap = True
        p2 = tx2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = BRAND_LIGHT
        p2.alignment = PP_ALIGN.CENTER

        if bullets:
            tx3 = slide.shapes.add_textbox(Inches(3), Inches(5.2), Inches(7), Inches(1.5)).text_frame
            tx3.word_wrap = True
            for b in bullets[:3]:
                pb = tx3.add_paragraph()
                pb.text = f"- {b}"
                pb.font.size = Pt(16)
                pb.font.color.rgb = WHITE
                pb.alignment = PP_ALIGN.CENTER

    def _add_text_only_slide(
        self,
        prs: Presentation,
        title: str,
        bullets: list[str],
        idx: int,
        total: int,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        self._add_accent_bar(slide)
        self._add_slide_title(slide, title)

        tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(4.0)).text_frame
        tx.word_wrap = True
        tx.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for b in bullets:
            pb = tx.add_paragraph()
            pb.text = f"- {b}"
            pb.font.size = Pt(20)
            pb.font.color.rgb = BODY_TEXT
            pb.space_after = Pt(12)

        self._add_slide_number(slide, idx, total)

    def _add_text_image_slide(
        self,
        prs: Presentation,
        title: str,
        bullets: list[str],
        image_path: str | None,
        idx: int,
        total: int,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        self._add_accent_bar(slide)
        self._add_slide_title(slide, title)

        # If no image, use full width for text
        if not image_path or not Path(image_path).exists():
            tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(4.0)).text_frame
            tx.word_wrap = True
            tx.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for b in bullets:
                pb = tx.add_paragraph()
                pb.text = f"- {b}"
                pb.font.size = Pt(20)
                pb.font.color.rgb = BODY_TEXT
                pb.space_after = Pt(12)
        else:
            # Text + Image layout
            tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(5.5), Inches(4.0)).text_frame
            tx.word_wrap = True
            tx.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for b in bullets:
                pb = tx.add_paragraph()
                pb.text = f"- {b}"
                pb.font.size = Pt(20)
                pb.font.color.rgb = BODY_TEXT
                pb.space_after = Pt(12)

            img_left = Inches(7.8)
            img_top = Inches(1.8)
            img_max_w = Inches(4.5)
            img_max_h = Inches(4.5)

            try:
                pic = slide.shapes.add_picture(image_path, img_left, img_top)
                ratio = min(img_max_w / float(pic.width), img_max_h / float(pic.height))
                new_w = int(pic.width * ratio)
                new_h = int(pic.height * ratio)
                pic.width = new_w
                pic.height = new_h
                pic.left = int(img_left + (img_max_w - new_w) / 2)
                pic.top = int(img_top + (img_max_h - new_h) / 2)
            except Exception as exc:
                logger.warning("Failed to embed image '%s': %s", image_path, exc)

        self._add_slide_number(slide, idx, total)

    def _add_image_placeholder(self, slide, query: str | None, left, top, width, height) -> None:
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = IMG_PLACEHOLDER_BG
        shape.line.fill.background()

        label = query or "image"
        tx = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.2), width - Inches(0.6), Inches(1.0)).text_frame
        tx.word_wrap = True
        p = tx.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = SUBTITLE_TEXT
        p.alignment = PP_ALIGN.CENTER

    def _add_accent_bar(self, slide) -> None:
        shape = slide.shapes.add_shape(1, Inches(1.5), Inches(0.9), Inches(0.6), Inches(0.08))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BRAND_PRIMARY
        shape.line.fill.background()

    def _add_slide_title(self, slide, title: str) -> None:
        tx = slide.shapes.add_textbox(Inches(1.5), Inches(1.1), Inches(10), Inches(0.8)).text_frame
        tx.word_wrap = True
        p = tx.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT

    def _add_highlight_box(self, slide, text: str, top_pos) -> None:
        shape = slide.shapes.add_shape(1, Inches(1.5), top_pos, Inches(10), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = HIGHLIGHT_BG
        shape.line.fill.background()

        tx = slide.shapes.add_textbox(Inches(1.8), top_pos + Inches(0.1), Inches(9.4), Inches(0.5)).text_frame
        tx.word_wrap = True
        p = tx.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = HIGHLIGHT_TEXT
        p.alignment = PP_ALIGN.LEFT

    def _add_slide_number(self, slide, idx: int, total: int) -> None:
        tx = slide.shapes.add_textbox(Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4)).text_frame
        p = tx.paragraphs[0]
        p.text = f"{idx + 1} / {total}"
        p.font.size = Pt(11)
        p.font.color.rgb = SUBTITLE_TEXT
        p.alignment = PP_ALIGN.RIGHT


def extract_image_queries(outline: dict) -> list[str]:
    queries: list[str] = []
    seen: set[str] = set()

    for slide in outline.get("slides", []):
        # New-format query field.
        q = (slide.get("image_query") or "").strip()
        if q and q not in seen:
            seen.add(q)
            queries.append(q)

        # Legacy elements fallback.
        for elem in slide.get("elements", []):
            if elem.get("type") in ("image_pexels", "image"):
                candidate = (elem.get("query") or "").strip()
                if candidate and candidate not in seen:
                    seen.add(candidate)
                    queries.append(candidate)

    return queries
